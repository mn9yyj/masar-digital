import streamlit as st
import requests, json, os, io, random
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- 1. واجهة المستخدم المتقدمة ---
st.set_page_config(page_title="المسار 🖥️ الرقمي - نسخة التصميم المتغير", layout="wide")
st.markdown("""
    <style>
    [data-testid="stSidebar"], footer, header {display: none !important;}
    .stApp { background: #0a0a0a; direction: rtl; }
    .brand { font-size: 45px; font-weight: 900; text-align: center; color: #48dbfb; margin-bottom: 20px; text-shadow: 2px 2px #222; }
    .centered-ui { max-width: 650px; margin: 0 auto; background: #1a1a1a; padding: 30px; border-radius: 20px; border: 1px solid #333; }
    .stButton>button { width: 100%; border-radius: 30px !important; background: linear-gradient(90deg, #1dd1a1, #10ac84) !important; height: 55px; font-weight: bold; border: none !important; font-size: 18px; }
    </style>
    """, unsafe_allow_html=True)

st.markdown('<div class="brand">المسار 🖥️ الرقمي - Pro Design</div>', unsafe_allow_html=True)

# --- 2. دالة بناء القوالب المتغيرة والترقيم ---
def apply_dynamic_layout(slide, item, index, lang_choice):
    # مصفوفة ألوان احترافية
    pro_colors = [RGBColor(72, 219, 251), RGBColor(255, 107, 107), RGBColor(29, 209, 161), RGBColor(254, 202, 87)]
    color = pro_colors[index % len(pro_colors)]
    
    # اختيار "نمط" الشريحة بناءً على رقمها لضمان التنوع
    style_type = index % 3 

    # أ) إضافة شكل هندسي متغير (رأس الشريحة)
    if style_type == 0: # نمط الدائرة الجانبية
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-0.5), Inches(-0.5), Inches(2.5), Inches(2.5))
    elif style_type == 1: # نمط الشريط العلوي
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
    else: # نمط الزاوية السفلية
        shape = slide.shapes.add_shape(MSO_SHAPE.TRIANGLE, Inches(8), Inches(5.5), Inches(2), Inches(2))
    
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.width = 0

    # ب) تنسيق العنوان
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(1))
    p_title = title_box.text_frame.paragraphs[0]
    p_title.text = str(item.get('title', 'تحليل البيانات'))
    p_title.font.size, p_title.font.bold = Pt(26), True
    p_title.alignment = PP_ALIGN.CENTER

    # ج) توزيع النقاط مع الترقيم (1, 2, 3...)
    points = item.get('points', [])
    for i, pt in enumerate(points[:4]):
        # تغيير مكان النص بناءً على الستايل
        left = Inches(0.8) if style_type != 2 else Inches(1.5)
        top = Inches(1.5 + (i * 1.3))
        
        # إضافة دائرة صغيرة للترقيم
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left - Inches(0.5), top, Inches(0.35), Inches(0.35))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = color
        
        # إضافة الرقم داخل الدائرة أو بجانبها
        num_box = slide.shapes.add_textbox(left - Inches(0.48), top - Inches(0.05), Inches(0.4), Inches(0.4))
        p_num = num_box.text_frame.paragraphs[0]
        p_num.text = str(i + 1)
        p_num.font.size, p_num.font.color.rgb = Pt(12), RGBColor(255, 255, 255)

        # إضافة نص النقطة (عربي ثم إنجليزي في المزدوج)
        box = slide.shapes.add_textbox(left, top, Inches(8.5), Inches(1))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        if lang_choice == "مزدوج":
            # دمج العربي فوق الإنجليزي في نقطة واحدة مرقمة
            ar_text = pt.get('ar', '')
            en_text = pt.get('en', '')
            p.text = f"{ar_text}\n{en_text}"
            p.font.size = Pt(12)
        else:
            p.text = str(pt)
            p.font.size = Pt(13)
        
        p.alignment = PP_ALIGN.RIGHT if lang_choice != "English" else PP_ALIGN.LEFT

# --- 3. محرك التشغيل ---
with st.container():
    st.markdown('<div class="centered-ui">', unsafe_allow_html=True)
    topic = st.text_input("🎯 موضوع العرض الرئيسي", placeholder="مثلاً: مكة المكرمة وتاريخها")
    slides_num = st.select_slider("عدد الشرائح المخصصة", options=[3, 5, 10, 15], value=5)
    lang = st.selectbox("🌐 لغة العرض", ["العربية", "English", "مزدوج"])
    
    if st.button("🚀 إصدار العرض المطور"):
        api_key = st.secrets.get("OPENROUTER_API_KEY")
        if not api_key: st.error("يرجى ضبط مفتاح API في الإعدادات")
        else:
            with st.spinner('🎨 جاري تنويع القوالب وترقيم النقاط...'):
                try:
                    # طلب بيانات منظمة جداً (نقاط مرقمة)
                    p_fmt = "{'ar': 'نص عربي', 'en': 'English text'}" if lang == "مزدوج" else "'نص النقطة'"
                    prompt = f"Create {slides_num} slides about '{topic}'. For each, provide 'title' and 4 'points' as a list of {p_fmt}. Return ONLY JSON."
                    
                    res = requests.post("https://openrouter.ai/api/v1/chat/completions",
                                        headers={"Authorization": f"Bearer {api_key}"},
                                        json={"model": "google/gemini-2.0-flash-001", "messages": [{"role": "user", "content": prompt}]})
                    
                    data = json.loads(res.json()['choices'][0]['message']['content'].replace("```json", "").replace("```", "").strip())
                    
                    prs = Presentation()
                    for i, slide_data in enumerate(data):
                        slide = prs.slides.add_slide(prs.slide_layouts[6])
                        apply_dynamic_layout(slide, slide_data, i, lang)
                    
                    buf = io.BytesIO()
                    prs.save(buf)
                    st.session_state['pro_file'] = buf.getvalue()
                except Exception as e: st.error(f"حدث خطأ في توليد البيانات: {e}")
    st.markdown('</div>', unsafe_allow_html=True)

if 'pro_file' in st.session_state:
    st.markdown('<div class="centered-ui" style="margin-top:20px;">', unsafe_allow_html=True)
    st.success("✅ تم الانتهاء! العرض الآن يحتوي على قوالب متنوعة ونقاط مرقمة.")
    st.download_button("📥 تحميل العرض (نسخة الإنفوجرافيك)", data=st.session_state['pro_file'], file_name="Digital_Path_Pro.pptx")
    st.markdown('</div>', unsafe_allow_html=True)
